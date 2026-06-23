# -*- coding: utf-8 -*-
"""PPTX export: Consulting-grade 17-slide GTM strategy presentation.

McKinsey/Bain/BCG/Gartner quality standard with visual frameworks, charts,
matrices, and minimal text per slide.
"""
import os
from typing import Any, Dict, List, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color palette - consulting grade
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x21, 0x29, 0x34)
MUTE = RGBColor(0x5A, 0x64, 0x72)
ACCENT = RGBColor(0x00, 0x7A, 0xFF)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0xF2, 0xF5, 0xFC)

W, H = Inches(13.333), Inches(7.5)


def _bg(slide, color):
    """Set slide background color."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add text box with multiple runs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    first = True
    for text, size, color, bold in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def _bullets(slide, x, y, w, h, items, size=14, color=INK, bullet="•", max_items=5):
    """Add bulleted list (max 5 items for consulting standard)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items[:max_items]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = f"{bullet} {it}"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _card(s, x, y, w, h, fill, border_color=None):
    """Add rounded rectangle card."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border_color:
        sh.line.color.rgb = border_color
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _kpi_card(s, x, y, w, h, value, label, color=ACCENT):
    """Add KPI card with large number and label."""
    _card(s, x, y, w, h, WHITE, MUTE)
    _text(s, x + Inches(0.15), y + Inches(0.2), w - Inches(0.3), Inches(0.8),
          [(value, 36, color, True)], align=PP_ALIGN.CENTER)
    _text(s, x + Inches(0.15), y + Inches(1.0), w - Inches(0.3), Inches(0.4),
          [(label, 12, MUTE, False)], align=PP_ALIGN.CENTER)


def _header(s, kicker, title):
    """Add slide header with kicker and title."""
    _bg(s, WHITE)
    _text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
          [(kicker.upper(), 13, MUTE, True)])
    _text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(0.8),
          [(title, 32, NAVY, True)])


def _content_slide(prs):
    """Create blank content slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_slide(prs, company, date_str):
    """Slide 1: Title slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, NAVY)
    _text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6),
          [("GO-TO-MARKET STRATEGY", 18, ICE, True)])
    _text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0),
          [(company, 48, WHITE, True)])
    _text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.8),
          [(date_str, 16, ICE, False)])


def _exec_summary(prs, r, gtm):
    """Slide 2: Executive Summary."""
    s = _content_slide(prs)
    _header(s, "Executive Summary", "Strategic Overview")
    
    # Key insights in cards
    _card(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(1.3), CARD_BG)
    _text(s, Inches(0.9), Inches(2.05), Inches(5.3), Inches(0.35),
          [("Market Opportunity", 14, NAVY, True)])
    summary_text = (r.get("executive_summary", "")[:200] + "...") if len(r.get("executive_summary", "")) > 200 else r.get("executive_summary", "—")
    _text(s, Inches(0.9), Inches(2.45), Inches(5.3), Inches(0.8),
          [(summary_text, 12, INK, False)])
    
    _card(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(1.3), CARD_BG)
    _text(s, Inches(7.1), Inches(2.05), Inches(5.3), Inches(0.35),
          [("Positioning", 14, NAVY, True)])
    pos = (gtm.get("foundation", {}) or {}).get("positioning_statement", "—")[:180]
    _text(s, Inches(7.1), Inches(2.45), Inches(5.3), Inches(0.8),
          [(pos, 12, INK, False)])
    
    # Key takeaway
    _text(s, Inches(0.7), Inches(3.5), Inches(11.9), Inches(0.4),
          [("Key Takeaway", 15, NAVY, True)])
    recs = r.get("recommendations", [])
    takeaway = recs[0] if recs else "Strategic GTM approach recommended"
    _text(s, Inches(0.7), Inches(4.0), Inches(11.9), Inches(2.0),
          [(takeaway, 16, INK, False)])


def _market_overview(prs, r):
    """Slide 3: Market Overview with KPI cards."""
    s = _content_slide(prs)
    _header(s, "Market Analysis", "Market Overview")
    
    # KPI cards - only 2 cards
    trends = r.get("market_trends", [])
    opps = r.get("opportunities", [])
    
    _kpi_card(s, Inches(2.5), Inches(2.0), Inches(3.5), Inches(1.5),
              f"{len(trends)}", "Key Trends", ACCENT)
    _kpi_card(s, Inches(6.8), Inches(2.0), Inches(3.5), Inches(1.5),
              f"{len(opps)}", "Opportunities", SUCCESS)
    
    # Key trends
    _text(s, Inches(0.7), Inches(3.8), Inches(11.9), Inches(0.4),
          [("Market Trends", 15, NAVY, True)])
    _bullets(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.5),
             trends[:4], size=13)


def _market_drivers(prs, r):
    """Slide 4: Market Drivers with visual framework."""
    s = _content_slide(prs)
    _header(s, "Market Dynamics", "Growth Drivers & Market Shifts")
    
    # Growth drivers
    _card(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(4.5), LIGHT_BG)
    _text(s, Inches(0.9), Inches(2.1), Inches(5.3), Inches(0.4),
          [("Growth Drivers", 16, NAVY, True)])
    opps = r.get("opportunities", [])
    _bullets(s, Inches(0.9), Inches(2.6), Inches(5.3), Inches(3.7),
             opps[:5], size=13)
    
    # Market shifts
    _card(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(4.5), LIGHT_BG)
    _text(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.4),
          [("Market Shifts", 16, NAVY, True)])
    risks = r.get("risks", [])
    _bullets(s, Inches(7.1), Inches(2.6), Inches(5.3), Inches(3.7),
             risks[:5], size=13)


def _icp_slide(prs, r, gtm):
    """Slide 5: Ideal Customer Profile."""
    s = _content_slide(prs)
    _header(s, "Target Customer", "Ideal Customer Profile")
    
    icp = (gtm.get("foundation", {}) or {}).get("icp", {}) or {}
    personas = r.get("buyer_personas", [])
    
    # ICP characteristics
    _card(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(2.0), CARD_BG)
    _text(s, Inches(0.9), Inches(2.05), Inches(5.3), Inches(0.35),
          [("Company Profile", 14, NAVY, True)])
    chars = [
        icp.get("primary_segment", "—"),
        icp.get("firmographics", "—")[:80],
        icp.get("technographics", "—")[:80]
    ]
    _bullets(s, Inches(0.9), Inches(2.5), Inches(5.3), Inches(1.4),
             [c for c in chars if c != "—"][:3], size=12)
    
    # Pain points
    _card(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(2.0), CARD_BG)
    _text(s, Inches(7.1), Inches(2.05), Inches(5.3), Inches(0.35),
          [("Key Pain Points", 14, NAVY, True)])
    pains = (gtm.get("foundation", {}) or {}).get("top_pains", [])
    _bullets(s, Inches(7.1), Inches(2.5), Inches(5.3), Inches(1.4),
             pains[:3], size=12)
    
    # Buying triggers
    _card(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.0), LIGHT_BG)
    _text(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.35),
          [("Buying Triggers", 14, NAVY, True)])
    triggers = (gtm.get("foundation", {}) or {}).get("trigger_events", [])
    _bullets(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.3),
             triggers[:4], size=12)


def _buying_committee(prs, gtm):
    """Slide 6: Buying Committee."""
    s = _content_slide(prs)
    _header(s, "Decision Makers", "Buying Committee")
    
    icp = (gtm.get("foundation", {}) or {}).get("icp", {}) or {}
    committee = icp.get("buying_committee", [])
    
    # Stakeholder cards
    roles = ["Economic Buyer", "Champion", "Influencer", "Decision Maker", "End User"]
    x_start = Inches(0.7)
    card_width = Inches(2.2)
    spacing = Inches(0.15)
    
    for i, role in enumerate(roles[:5]):
        x = x_start + i * (card_width + spacing)
        _card(s, x, Inches(2.0), card_width, Inches(3.5), CARD_BG)
        _text(s, x + Inches(0.1), Inches(2.15), card_width - Inches(0.2), Inches(0.5),
              [(role, 13, NAVY, True)], align=PP_ALIGN.CENTER)
        
        # Add committee member if available
        member = committee[i] if i < len(committee) else "TBD"
        _text(s, x + Inches(0.1), Inches(2.7), card_width - Inches(0.2), Inches(2.6),
              [(member[:60], 11, INK, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # Key takeaway
    _text(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(0.8),
          [("Multi-threaded approach required across technical and business stakeholders", 13, MUTE, False)])


def _competitive_landscape(prs, r):
    """Slide 7: Competitive Landscape Matrix."""
    s = _content_slide(prs)
    _header(s, "Competition", "Competitive Landscape")
    
    comps = (r.get("company_competitors", []) or []) + (r.get("product_competitors", []) or [])
    
    if comps:
        # Competitor matrix
        rows = [["Competitor", "Positioning", "Strength", "Weakness"]]
        for c in comps[:5]:
            swot = c.get("swot", {}) or {}
            strengths = swot.get("strengths", [])
            weaknesses = swot.get("weaknesses", [])
            
            strength = strengths[0].get("point", "—") if strengths else "—"
            weakness = weaknesses[0].get("point", "—") if weaknesses else "—"
            
            rows.append([
                c.get("name", "—")[:25],
                c.get("value_proposition", "—")[:40],
                strength[:35],
                weakness[:35]
            ])
        
        tbl = s.shapes.add_table(len(rows), 4, Inches(0.7), Inches(1.9),
                                 Inches(11.9), Inches(0.45) * len(rows)).table
        tbl.columns[0].width = Inches(2.5)
        tbl.columns[1].width = Inches(3.5)
        tbl.columns[2].width = Inches(3.0)
        tbl.columns[3].width = Inches(2.9)
        
        # Header row
        for ci, head in enumerate(rows[0]):
            cell = tbl.cell(0, ci)
            cell.text = head
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        
        # Data rows
        for ri in range(1, len(rows)):
            for ci in range(4):
                cell = tbl.cell(ri, ci)
                cell.text = rows[ri][ci]
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = INK


def _opportunity_gaps(prs, r, gtm):
    """Slide 8: Opportunity Gaps."""
    s = _content_slide(prs)
    _header(s, "Strategic Opportunities", "White Space & Market Gaps")
    
    # Opportunities
    _card(s, Inches(0.7), Inches(1.9), Inches(5.7), Inches(4.5), LIGHT_BG)
    _text(s, Inches(0.9), Inches(2.1), Inches(5.3), Inches(0.4),
          [("Market Opportunities", 16, SUCCESS, True)])
    opps = r.get("opportunities", [])
    _bullets(s, Inches(0.9), Inches(2.6), Inches(5.3), Inches(3.7),
             opps[:5], size=13, color=INK)
    
    # Competitive gaps
    _card(s, Inches(6.9), Inches(1.9), Inches(5.7), Inches(4.5), LIGHT_BG)
    _text(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(0.4),
          [("Competitive Weaknesses", 16, WARNING, True)])
    
    # Extract competitor weaknesses
    comps = (r.get("company_competitors", []) or []) + (r.get("product_competitors", []) or [])
    weaknesses = []
    for c in comps[:3]:
        swot = c.get("swot", {}) or {}
        for w in (swot.get("weaknesses", []) or [])[:2]:
            point = w.get("point", "") if isinstance(w, dict) else str(w)
            if point:
                weaknesses.append(f"{c.get('name', 'Competitor')}: {point[:50]}")
    
    _bullets(s, Inches(7.1), Inches(2.6), Inches(5.3), Inches(3.7),
             weaknesses[:5] or ["Analyze competitor gaps"], size=12, color=INK)


def _positioning_statement(prs, gtm):
    """Slide 9: Positioning Statement."""
    s = _content_slide(prs)
    _header(s, "Positioning", "Strategic Positioning")
    
    foundation = gtm.get("foundation", {}) or {}
    slot = foundation.get("slot_statement", {}) or {}
    
    # Positioning framework
    _card(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.8), LIGHT_BG)
    
    y = Inches(2.5)
    elements = [
        ("For", slot.get("for_who", "—")),
        ("Who need", slot.get("who_need", "—")),
        ("Our solution", slot.get("category", "—")),
        ("That delivers", slot.get("promise", "—")),
        ("Unlike", slot.get("unlike", "—")),
        ("We provide", slot.get("proof", "—"))
    ]
    
    for label, value in elements:
        _text(s, Inches(1.8), y, Inches(2.0), Inches(0.4),
              [(label, 13, NAVY, True)])
        _text(s, Inches(4.0), y, Inches(7.5), Inches(0.4),
              [(value[:80], 12, INK, False)])
        y += Inches(0.55)


def _messaging_architecture(prs, gtm):
    """Slide 10: Messaging Architecture."""
    s = _content_slide(prs)
    _header(s, "Messaging", "Messaging Architecture")
    
    activation = gtm.get("activation", {}) or {}
    messaging = activation.get("messaging_by_persona", [])
    
    if messaging:
        persona_msg = messaging[0]
        
        # Core promise (top of pyramid)
        _card(s, Inches(4.5), Inches(2.0), Inches(4.3), Inches(0.8), ACCENT)
        _text(s, Inches(4.6), Inches(2.15), Inches(4.1), Inches(0.5),
              [(persona_msg.get("core_promise", "Core Promise")[:60], 14, WHITE, True)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
        # Supporting pillars
        pillars = persona_msg.get("pillars", [])
        x_start = Inches(2.5)
        for i, pillar in enumerate(pillars[:3]):
            x = x_start + i * Inches(3.0)
            _card(s, x, Inches(3.2), Inches(2.7), Inches(1.0), CARD_BG)
            _text(s, x + Inches(0.1), x + Inches(3.4), Inches(2.5), Inches(0.6),
                  [(pillar[:50], 11, INK, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
        # Proof points
        _text(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(0.4),
              [("Proof Points", 15, NAVY, True)])
        proofs = persona_msg.get("proof_points", [])
        _bullets(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.5),
                 proofs[:4], size=12)


def _strategic_priorities(prs, gtm):
    """Slide 11: Strategic Priorities."""
    s = _content_slide(prs)
    _header(s, "Strategy", "Strategic Priorities")
    
    foundation = gtm.get("foundation", {}) or {}
    beachhead = foundation.get("beachhead", {}) or {}
    
    priorities = [
        f"1. {beachhead.get('segment', 'Target beachhead segment')}",
        f"2. {beachhead.get('entry_wedge', 'Execute entry strategy')}",
        "3. Build competitive moat through differentiation",
        "4. Scale through proven channels",
        "5. Expand to adjacent segments"
    ]
    
    y = Inches(2.2)
    for i, priority in enumerate(priorities[:5]):
        color = SUCCESS if i == 0 else ACCENT if i == 1 else MUTE
        _card(s, Inches(0.7), y, Inches(11.9), Inches(0.8), LIGHT_BG)
        _text(s, Inches(1.0), y + Inches(0.15), Inches(11.3), Inches(0.5),
              [(priority, 15, color, i < 2)])
        y += Inches(1.0)


def _channel_strategy(prs, gtm):
    """Slide 9: Channel Strategy - Enhanced design."""
    s = _content_slide(prs)
    _header(s, "Go-To-Market", "Channel Strategy & Mix")
    
    activation = gtm.get("activation", {}) or {}
    channels = activation.get("channel_plays", [])
    
    # Primary channel (large card)
    if channels:
        primary = channels[0]
        _card(s, Inches(0.7), Inches(2.0), Inches(5.5), Inches(4.2), ACCENT)
        _text(s, Inches(0.9), Inches(2.3), Inches(5.1), Inches(0.5),
              [("PRIMARY CHANNEL", 14, WHITE, True)])
        _text(s, Inches(0.9), Inches(2.9), Inches(5.1), Inches(0.7),
              [(primary.get("channel", "Channel"), 24, WHITE, True)])
        _text(s, Inches(0.9), Inches(3.8), Inches(5.1), Inches(0.4),
              [("Funnel Role", 12, ICE, True)])
        _text(s, Inches(0.9), Inches(4.2), Inches(5.1), Inches(0.9),
              [(primary.get("funnel_role", "—")[:80], 13, WHITE, False)])
        _text(s, Inches(0.9), Inches(5.2), Inches(5.1), Inches(0.4),
              [("Investment Focus", 12, ICE, True)])
        _text(s, Inches(0.9), Inches(5.6), Inches(5.1), Inches(0.5),
              [(primary.get("invest", "—")[:80], 13, WHITE, False)])
    
    # Secondary channels (stacked cards)
    y = Inches(2.0)
    for i, ch in enumerate(channels[1:4]):
        _card(s, Inches(6.7), y, Inches(5.9), Inches(1.3), LIGHT_BG)
        _text(s, Inches(6.9), y + Inches(0.15), Inches(5.5), Inches(0.4),
              [(ch.get("channel", "Channel"), 15, NAVY, True)])
        _text(s, Inches(6.9), y + Inches(0.6), Inches(5.5), Inches(0.5),
              [(ch.get("funnel_role", "—")[:70], 12, INK, False)])
        y += Inches(1.4)
    
    # Key takeaway
    _text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.8),
          [("Multi-channel approach with clear primary focus and supporting channels", 13, MUTE, False)])


def _campaign_recommendations(prs, gtm):
    """Slide 13: Campaign Recommendations."""
    s = _content_slide(prs)
    _header(s, "Campaigns", "Campaign Recommendations")
    
    execution = gtm.get("execution", {}) or {}
    demand = execution.get("demand_gen", {}) or {}
    campaigns = demand.get("campaign_concepts", [])
    
    # Campaign cards
    y = Inches(2.0)
    for i, campaign in enumerate(campaigns[:4]):
        _card(s, Inches(0.7), y, Inches(11.9), Inches(1.0), LIGHT_BG)
        _text(s, Inches(1.0), y + Inches(0.15), Inches(0.8), Inches(0.7),
              [(f"#{i+1}", 24, ACCENT, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(2.0), y + Inches(0.2), Inches(10.3), Inches(0.6),
              [(campaign[:100], 13, INK, False)])
        y += Inches(1.15)


def _sales_motion(prs, gtm):
    """Slide 14: Sales Motion."""
    s = _content_slide(prs)
    _header(s, "Sales", "Sales Motion & Process")
    
    execution = gtm.get("execution", {}) or {}
    playbook = execution.get("sales_playbook", {}) or {}
    stages = playbook.get("stages", [])
    
    # Process flow
    x = Inches(0.7)
    card_width = Inches(2.1)
    for i, stage in enumerate(stages[:5]):
        color = SUCCESS if i == 0 else ACCENT if i < 3 else MUTE
        _card(s, x, Inches(2.2), card_width, Inches(1.5), LIGHT_BG, color)
        _text(s, x + Inches(0.1), Inches(2.35), card_width - Inches(0.2), Inches(0.4),
              [(stage.get("stage", f"Stage {i+1}"), 12, color, True)],
              align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.1), Inches(2.8), card_width - Inches(0.2), Inches(0.8),
              [(stage.get("objective", "—")[:50], 10, INK, False)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        
        # Arrow
        if i < len(stages[:5]) - 1:
            _text(s, x + card_width, Inches(2.7), Inches(0.3), Inches(0.4),
                  [("→", 20, MUTE, False)], align=PP_ALIGN.CENTER)
        x += card_width + Inches(0.3)
    
    # Qualification framework
    _text(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.4),
          [("Qualification Framework", 15, NAVY, True)])
    _text(s, Inches(0.7), Inches(4.7), Inches(11.9), Inches(1.5),
          [(playbook.get("qualification_framework", "—"), 13, INK, False)])


def _roadmap_90day(prs, gtm):
    """Slide 15: 90-Day GTM Roadmap."""
    s = _content_slide(prs)
    _header(s, "Execution", "90-Day GTM Roadmap")
    
    execution = gtm.get("execution", {}) or {}
    roadmap = execution.get("roadmap_90day", [])
    
    # Timeline view
    x = Inches(0.7)
    card_width = Inches(3.8)
    colors = [SUCCESS, ACCENT, WARNING]
    
    for i, phase in enumerate(roadmap[:3]):
        color = colors[i] if i < len(colors) else MUTE
        _card(s, x, Inches(2.0), card_width, Inches(4.0), LIGHT_BG, color)
        _text(s, x + Inches(0.2), Inches(2.2), card_width - Inches(0.4), Inches(0.5),
              [(phase.get("phase", f"Month {i+1}"), 16, color, True)])
        _text(s, x + Inches(0.2), Inches(2.8), card_width - Inches(0.4), Inches(0.6),
              [(phase.get("objective", "—")[:80], 12, INK, False)])
        
        # Workstreams
        workstreams = phase.get("workstreams", [])
        _bullets(s, x + Inches(0.2), Inches(3.6), card_width - Inches(0.4), Inches(2.2),
                 [w.get("workstream", "")[:40] for w in workstreams[:4]], size=11)
        x += card_width + Inches(0.2)


def _kpis_metrics(prs, gtm):
    """Slide 16: KPIs & Success Metrics."""
    s = _content_slide(prs)
    _header(s, "Metrics", "KPIs & Success Metrics")
    
    execution = gtm.get("execution", {}) or {}
    metrics = execution.get("metrics", {}) or {}
    
    # North star
    _card(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(1.0), ACCENT)
    _text(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(0.6),
          [(f"North Star: {metrics.get('north_star_metric', 'Revenue Growth')}", 18, WHITE, True)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    
    # KPI cards
    input_metrics = metrics.get("input_metrics", [])
    funnel_kpis = metrics.get("funnel_kpis", [])
    
    all_metrics = input_metrics + funnel_kpis
    x = Inches(0.7)
    card_width = Inches(2.8)
    
    for i, metric in enumerate(all_metrics[:4]):
        _card(s, x, Inches(3.3), card_width, Inches(2.5), CARD_BG)
        _text(s, x + Inches(0.15), Inches(3.5), card_width - Inches(0.3), Inches(0.5),
              [(metric.get("metric", "Metric"), 13, NAVY, True)], align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.15), Inches(4.1), card_width - Inches(0.3), Inches(0.6),
              [(metric.get("target_band", "—"), 16, ACCENT, True)], align=PP_ALIGN.CENTER)
        _text(s, x + Inches(0.15), Inches(4.8), card_width - Inches(0.3), Inches(0.8),
              [(metric.get("why", "")[:60], 10, MUTE, False)], align=PP_ALIGN.CENTER)
        x += card_width + Inches(0.15)


def _sources_methodology(prs, result):
    """Slide 17: Sources & Methodology."""
    s = _content_slide(prs)
    _header(s, "Appendix", "Sources & Methodology")
    
    sources = result.get("sources", [])
    
    _text(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.4),
          [("Research Methodology", 15, NAVY, True)])
    _text(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(1.0),
          [("This analysis is based on comprehensive market research including competitive intelligence, "
           "customer insights, market trends, and strategic frameworks.", 13, INK, False)])
    
    _text(s, Inches(0.7), Inches(3.6), Inches(11.9), Inches(0.4),
          [("Key Sources", 15, NAVY, True)])
    
    # Source list
    source_list = []
    for i, src in enumerate(sources[:15]):
        title = src.get("title", "")
        domain = src.get("domain", "")
        if title and domain:
            source_list.append(f"{domain}: {title[:60]}")
        elif domain:
            source_list.append(domain)
    
    if source_list:
        _bullets(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.5),
                 source_list[:10], size=11, max_items=10)
    else:
        _text(s, Inches(0.7), Inches(4.1), Inches(11.9), Inches(1.0),
              [("Primary and secondary research sources", 12, MUTE, False)])


def export_pptx(result: Dict[str, Any], path: Optional[str] = None) -> str:
    """Export consulting-grade 13-slide GTM strategy presentation."""
    r = result.get("report", {}) or {}
    plan = result.get("plan", {}) or {}
    gtm = result.get("gtm_strategy") or {}
    
    company = plan.get("subject_entity") or r.get("title") or "Market Analysis"
    date_str = datetime.now().strftime("%B %Y")
    
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    
    # Generate 13 slides (removed: Buying Committee, Competitive Landscape, Messaging Architecture, Strategic Priorities)
    _title_slide(prs, company, date_str)                    # 1
    _exec_summary(prs, r, gtm)                              # 2
    _market_overview(prs, r)                                # 3
    _market_drivers(prs, r)                                 # 4
    _icp_slide(prs, r, gtm)                                 # 5
    # _buying_committee removed                             # (deleted)
    # _competitive_landscape removed                        # (deleted)
    _opportunity_gaps(prs, r, gtm)                          # 6
    _positioning_statement(prs, gtm)                        # 7
    # _messaging_architecture removed                       # (deleted)
    # _strategic_priorities removed                         # (deleted)
    _channel_strategy(prs, gtm)                             # 8 (redesigned)
    _campaign_recommendations(prs, gtm)                     # 9
    _sales_motion(prs, gtm)                                 # 10
    _roadmap_90day(prs, gtm)                                # 11
    _kpis_metrics(prs, gtm)                                 # 12
    _sources_methodology(prs, result)                       # 13
    
    out = path or os.path.join("exports", "pptx",
                               f"gtm_strategy_{company.lower().replace(' ', '_')[:40]}.pptx")
    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    prs.save(out)
    print(f"   ✓ Consulting-grade GTM deck (13 slides): {out}")
    return out

# Made with Bob
